from io import BytesIO
from typing import Optional, Union
import pandas as pd
from pptx import Presentation
from pptx.slide import Slide
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

from src.powerpoint.constants.layouts import SlideLayouts
from src.powerpoint.constants.sizes import STANDARD_CHART_SIZES
from src.powerpoint.numpy.bridge import convert_series_to_list


class ThemedPresentation:
    def __init__(self, pptx: Optional[str] = None, *args, **kwargs):
        self.prs = Presentation(pptx)

    def add_bar_chart(self) -> None:
        raise NotImplementedError()

    def add_line_chart(
        self,
        data: pd.DataFrame,
        x_axis: str,
        y_value: str,
        hue: Optional[str] = None,
        title: Optional[str] = None,
        legend: Optional[bool] = True,
        y_axis_reference_value: Optional[Union[int, float]] = None,
        y_axis_reference_name: Optional[str] = None,
    ) -> None:
        assert (
            x_axis in data.columns
        ), f"Requested x-axis of {x_axis} not found in source columns {data.columns}"
        assert (
            y_value in data.columns
        ), f"Requested y-value of {y_value} not found in source columns {data.columns}"

        slide = self.add_slide(SlideLayouts.TITLE_ONLY)
        if title:
            title = slide.shapes.title
            title.text = title

        chart_data = ChartData()
        data.sort_values(by=x_axis, inplace=True)

        categories = data[x_axis].unique()
        chart_data.categories = convert_series_to_list(categories)

        if hue:
            assert (
                hue in data.columns
            ), f"Requested hue of {hue} not found in source columns {data.columns}"

            hues = data[hue].unique()
            for hue in hues:
                data_to_add = list(data.loc[data[hue] == hue, y_value])
                chart_data.add_series("Total", data_to_add)
        else:
            dat = list(data[y_value])
            chart_data.add_series("Total", dat)

        if y_axis_reference_value:
            dat = [y_axis_reference_value for _ in range(len(categories))]
            chart_data.add_series(y_axis_reference_name, dat)

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE, *STANDARD_CHART_SIZES, chart_data
        ).chart

        if legend:
            chart.has_legend = True
            chart.legend.include_in_layout = False

    def add_slide(self, layout: SlideLayouts) -> Slide:
        return self.prs.slides.add_slide(self.prs.slide_layouts[layout])

    def save(self, destination: Union[str, BytesIO]) -> None:
        self.prs.save(destination)

    def to_buffer(self) -> BytesIO:
        buffer = BytesIO()
        self.prs.save(buffer)
        return buffer
